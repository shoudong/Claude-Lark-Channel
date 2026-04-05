#!/usr/bin/env python3
"""
Generate Claude Channel Setup Guide presentation for ADVANCE.AI leadership.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# === Constants ===
OUTPUT_PATH = "/Users/dong-ai/Documents/ADVANCE-AI-OS/09_Speeches_Writing/Claude-Channel-Setup-Guide.pptx"

# Colors
DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
MEDIUM_BLUE = RGBColor(0x2B, 0x6C, 0xB0)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
ACCENT_GREEN = RGBColor(0x38, 0xA1, 0x69)
ACCENT_ORANGE = RGBColor(0xDD, 0x6B, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x71, 0x71, 0x7A)
LIGHT_GRAY = RGBColor(0xF7, 0xFA, 0xFC)
TABLE_HEADER_BG = DARK_BLUE
TABLE_ALT_ROW = RGBColor(0xEB, 0xF4, 0xFF)
BORDER_COLOR = RGBColor(0xCB, 0xD5, 0xE0)

# Fonts
FONT_FAMILY = "Calibri"
FONT_HEADING = "Calibri"

# Slide dimensions (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # === Slide 1: Title ===
    slide1 = prs.slides.add_slide(blank_layout)
    add_title_slide(slide1)

    # === Slide 2: What is Claude Channel? ===
    slide2 = prs.slides.add_slide(blank_layout)
    add_content_slide(
        slide2,
        "What is Claude Channel?",
        [
            "A dedicated Lark group chat where you message Claude directly",
            "Claude processes your instruction and replies in seconds",
            "Simple tasks (calendar, tasks, email) -- instant, zero AI cost",
            "Complex tasks (analysis, drafting, summarization) -- full AI power",
            "Works from any device with Lark (phone, desktop, web)",
            "Replaces unreliable automation tools (OpenClaw)",
        ],
    )

    # === Slide 3: Architecture Overview ===
    slide3 = prs.slides.add_slide(blank_layout)
    add_architecture_slide(slide3)

    # === Slide 4: Simple Tasks ===
    slide4 = prs.slides.add_slide(blank_layout)
    add_simple_tasks_slide(slide4)

    # === Slide 5: Complex Tasks ===
    slide5 = prs.slides.add_slide(blank_layout)
    add_complex_tasks_slide(slide5)

    # === Slide 6: Token Efficiency ===
    slide6 = prs.slides.add_slide(blank_layout)
    add_token_efficiency_slide(slide6)

    # === Slide 7: Security ===
    slide7 = prs.slides.add_slide(blank_layout)
    add_security_slide(slide7)

    # === Slide 8: Setup Steps 1/2 ===
    slide8 = prs.slides.add_slide(blank_layout)
    add_setup_steps_1_slide(slide8)

    # === Slide 9: Setup Steps 2/2 ===
    slide9 = prs.slides.add_slide(blank_layout)
    add_setup_steps_2_slide(slide9)

    # === Slide 10: Architecture Evolution ===
    slide10 = prs.slides.add_slide(blank_layout)
    add_evolution_slide(slide10)

    # === Slide 11: What's Next ===
    slide11 = prs.slides.add_slide(blank_layout)
    add_whats_next_slide(slide11)

    # === Slide 12: Thank You ===
    slide12 = prs.slides.add_slide(blank_layout)
    add_thank_you_slide(slide12)

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")


# === Helper Functions ===

def set_font(run, size=18, bold=False, italic=False, color=BLACK, font_family=FONT_FAMILY):
    """Apply font properties to a run."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_family


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=BLACK,
                alignment=PP_ALIGN.LEFT, font_family=FONT_FAMILY, italic=False):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font_family=font_family, italic=italic)
    return txBox


def add_header_bar(slide, title_text):
    """Add the dark blue header bar with title."""
    # Full-width header bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title text in the bar
    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.2),
        Inches(11), Inches(0.7),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_font(run, size=28, bold=True, color=WHITE, font_family=FONT_HEADING)


def add_footer_bar(slide):
    """Add a subtle footer with page context."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.1),
        SLIDE_WIDTH, Inches(0.4),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_GRAY
    bar.line.fill.background()

    add_textbox(
        slide, Inches(0.7), Inches(7.15),
        Inches(5), Inches(0.3),
        "ADVANCE.AI  |  Claude Channel Setup Guide  |  April 2026",
        size=9, color=GRAY,
    )


def add_bullet_slide(slide, title, bullets, start_top=Inches(1.4), left=Inches(0.9),
                      width=Inches(11.5), bullet_size=17, spacing=Pt(8)):
    """Add header + bullet points."""
    add_header_bar(slide, title)
    add_footer_bar(slide)

    txBox = slide.shapes.add_textbox(left, start_top, width, Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = spacing
        p.level = 0

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u2022  "
        set_font(run_bullet, size=bullet_size, color=MEDIUM_BLUE, bold=True)

        run_text = p.add_run()
        run_text.text = bullet
        set_font(run_text, size=bullet_size, color=BLACK)


def add_content_slide(slide, title, bullets):
    """Standard content slide with header bar and bullets."""
    add_bullet_slide(slide, title, bullets)


def add_rounded_rect(slide, left, top, width, height, fill_color, text,
                     font_size=13, font_color=WHITE, bold=True):
    """Add a rounded rectangle shape with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.line.width = Pt(1)

    # Adjust corner radius
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Vertical centering
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)

    run = tf.paragraphs[0].add_run()
    run.text = text
    set_font(run, size=font_size, bold=bold, color=font_color)

    # Try to vertically center
    tf_element = tf._txBody
    bodyPr = tf_element.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top, color=MEDIUM_BLUE, width=Pt(2.5)):
    """Add an arrow connector between two points."""
    # Use a line shape with arrow
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        start_left, start_top,
        end_left - start_left, Pt(20),
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def add_down_arrow(slide, left, top, width=Pt(20), height=Inches(0.35), color=MEDIUM_BLUE):
    """Add a downward arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# === Slide Builders ===

def add_title_slide(slide):
    """Slide 1: Title slide."""
    # Full background dark blue
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(2.6),
        Inches(1.5), Pt(4),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GREEN
    accent.line.fill.background()

    # Title
    add_textbox(
        slide, Inches(1.2), Inches(2.9),
        Inches(10), Inches(1.2),
        "Claude Channel",
        size=48, bold=True, color=WHITE, font_family=FONT_HEADING,
    )

    # Dash + subtitle
    add_textbox(
        slide, Inches(1.2), Inches(3.9),
        Inches(10), Inches(0.7),
        "AI Assistant in Lark",
        size=30, bold=False, color=RGBColor(0xA0, 0xC4, 0xE8), font_family=FONT_HEADING,
    )

    # Subtitle description
    add_textbox(
        slide, Inches(1.2), Inches(4.7),
        Inches(10), Inches(0.5),
        "Setup Guide for ADVANCE.AI Leadership Team",
        size=18, color=RGBColor(0xCB, 0xD5, 0xE0),
    )

    # Date and author
    add_textbox(
        slide, Inches(1.2), Inches(5.6),
        Inches(4), Inches(0.4),
        "April 3, 2026",
        size=14, color=RGBColor(0x90, 0xA0, 0xB0),
    )

    add_textbox(
        slide, Inches(1.2), Inches(6.0),
        Inches(4), Inches(0.4),
        "Dong Shou",
        size=14, color=RGBColor(0x90, 0xA0, 0xB0), bold=True,
    )


def add_architecture_slide(slide):
    """Slide 3: Architecture Overview with flow diagram."""
    add_header_bar(slide, "Architecture Overview")
    add_footer_bar(slide)

    # Layout dimensions for the flow
    box_h = Inches(0.65)
    arrow_w = Inches(0.7)
    row1_top = Inches(2.4)
    row2_top = Inches(4.5)

    # --- Row 1: Main flow ---
    # Box: You / Lark
    x1 = Inches(0.8)
    bw1 = Inches(2.0)
    add_rounded_rect(slide, x1, row1_top, bw1, box_h, MEDIUM_BLUE, "You / Lark", font_size=14)

    # Arrow 1
    a1_left = x1 + bw1 + Inches(0.15)
    add_arrow(slide, a1_left, row1_top + Inches(0.22), a1_left + arrow_w, row1_top + Inches(0.22))

    # Box: ngrok Tunnel
    x2 = a1_left + arrow_w + Inches(0.15)
    bw2 = Inches(2.2)
    add_rounded_rect(slide, x2, row1_top, bw2, box_h, RGBColor(0x55, 0x55, 0x80), "ngrok Tunnel (TLS)", font_size=13)

    # Arrow 2
    a2_left = x2 + bw2 + Inches(0.15)
    add_arrow(slide, a2_left, row1_top + Inches(0.22), a2_left + arrow_w, row1_top + Inches(0.22))

    # Box: Channel Server
    x3 = a2_left + arrow_w + Inches(0.15)
    bw3 = Inches(2.5)
    add_rounded_rect(slide, x3, row1_top, bw3, box_h, DARK_BLUE, "Channel Server", font_size=14)

    # Arrow 3
    a3_left = x3 + bw3 + Inches(0.15)
    add_arrow(slide, a3_left, row1_top + Inches(0.22), a3_left + arrow_w, row1_top + Inches(0.22))

    # Box: Claude Code Session
    x4 = a3_left + arrow_w + Inches(0.15)
    bw4 = Inches(2.8)
    add_rounded_rect(slide, x4, row1_top, bw4, box_h, ACCENT_GREEN, "Claude Code Session", font_size=14)

    # --- Row 2: Branching connections ---
    # Down arrow from Channel Server to lark-cli
    cs_center_x = x3 + bw3 / 2 - Pt(10)
    add_down_arrow(slide, cs_center_x, row1_top + box_h + Inches(0.1), Pt(20), Inches(0.5))

    # Box: lark-cli (Simple Tasks)
    lark_cli_w = Inches(2.5)
    lark_cli_x = x3
    add_rounded_rect(slide, lark_cli_x, row2_top, lark_cli_w, box_h,
                     ACCENT_ORANGE, "lark-cli (Simple Tasks)", font_size=13)

    # Down arrow from Claude Code to Obsidian
    cc_center_x = x4 + bw4 / 2 - Pt(10)
    add_down_arrow(slide, cc_center_x, row1_top + box_h + Inches(0.1), Pt(20), Inches(0.5))

    # Box: Obsidian Vault
    obs_w = Inches(2.8)
    obs_x = x4
    add_rounded_rect(slide, obs_x, row2_top, obs_w, box_h,
                     RGBColor(0x6B, 0x46, 0xC1), "Obsidian Vault (Memory)", font_size=13)

    # Tagline at bottom
    add_textbox(
        slide, Inches(0.8), Inches(5.8),
        Inches(11), Inches(0.5),
        "3 components.  1 process.  Zero idle cost.",
        size=20, bold=True, color=DARK_BLUE, font_family=FONT_HEADING,
        alignment=PP_ALIGN.CENTER,
    )

    # Labels for the arrows
    add_textbox(slide, a1_left, row1_top - Inches(0.35), arrow_w, Inches(0.3),
                "webhook", size=9, color=GRAY, alignment=PP_ALIGN.CENTER, italic=True)
    add_textbox(slide, a2_left, row1_top - Inches(0.35), arrow_w, Inches(0.3),
                "localhost", size=9, color=GRAY, alignment=PP_ALIGN.CENTER, italic=True)
    add_textbox(slide, a3_left, row1_top - Inches(0.35), arrow_w, Inches(0.3),
                "stdin/out", size=9, color=GRAY, alignment=PP_ALIGN.CENTER, italic=True)


def add_simple_tasks_slide(slide):
    """Slide 4: How It Works -- Simple Tasks."""
    add_header_bar(slide, "How It Works -- Simple Tasks")
    add_footer_bar(slide)

    # Left column: explanation
    bullets = [
        "Keywords detected: calendar, tasks, email, agenda, schedule, todo",
        "Handled directly by the channel server via lark-cli",
        "No AI tokens consumed",
        "Response time: ~1.5 seconds",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(6.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        rb = p.add_run()
        rb.text = "\u2022  "
        set_font(rb, size=16, color=MEDIUM_BLUE, bold=True)
        rt = p.add_run()
        rt.text = bullet
        set_font(rt, size=16, color=BLACK)

    # Right column: examples in styled boxes
    example_top = Inches(1.5)
    ex_left = Inches(8.2)
    ex_width = Inches(4.5)

    # Examples header
    add_textbox(slide, ex_left, Inches(1.3), ex_width, Inches(0.4),
                "Examples", size=18, bold=True, color=DARK_BLUE)

    examples = [
        ('"show my calendar"', "instant agenda"),
        ('"check my tasks"', "task list with overdue flags"),
        ('"show my inbox"', "latest 5 emails"),
    ]

    for i, (cmd, result) in enumerate(examples):
        y = Inches(1.9) + i * Inches(1.05)

        # Command box
        cmd_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ex_left, y, ex_width, Inches(0.4),
        )
        cmd_box.fill.solid()
        cmd_box.fill.fore_color.rgb = RGBColor(0x2D, 0x37, 0x48)
        cmd_box.line.fill.background()
        cmd_box.adjustments[0] = 0.15
        cmd_tf = cmd_box.text_frame
        cmd_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        cmd_run = cmd_tf.paragraphs[0].add_run()
        cmd_run.text = "  " + cmd
        set_font(cmd_run, size=13, color=ACCENT_GREEN, font_family="Courier New")
        bodyPr = cmd_tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
            bodyPr.set('lIns', str(Emu(Inches(0.15))))

        # Result label
        add_textbox(slide, ex_left + Inches(0.2), y + Inches(0.45), ex_width, Inches(0.3),
                    "\u2192  " + result, size=12, color=GRAY, italic=True)

    # Green accent badge
    add_rounded_rect(slide, Inches(0.9), Inches(4.5), Inches(2.5), Inches(0.5),
                     ACCENT_GREEN, "0 tokens  |  Free", font_size=15, font_color=WHITE)


def add_complex_tasks_slide(slide):
    """Slide 5: How It Works -- Complex Tasks."""
    add_header_bar(slide, "How It Works -- Complex Tasks")
    add_footer_bar(slide)

    bullets = [
        "Anything not matching simple keywords \u2192 pushed to Claude Code",
        "Claude processes with full tool access (files, search, lark-cli, web)",
        "Runs in a warm session -- no cold start overhead",
        "Response time: 5-30 seconds depending on complexity",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(6.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        rb = p.add_run()
        rb.text = "\u2022  "
        set_font(rb, size=16, color=MEDIUM_BLUE, bold=True)
        rt = p.add_run()
        rt.text = bullet
        set_font(rt, size=16, color=BLACK)

    # Right column: examples
    ex_left = Inches(8.2)
    ex_width = Inches(4.5)
    add_textbox(slide, ex_left, Inches(1.3), ex_width, Inches(0.4),
                "Examples", size=18, bold=True, color=DARK_BLUE)

    examples = [
        '"summarize my week from March 30 to today and save"',
        '"draft a reply to Sean about the Q2 pipeline"',
        '"analyze the Indonesia MMR data and note key takeaways"',
    ]

    for i, cmd in enumerate(examples):
        y = Inches(1.9) + i * Inches(0.9)
        cmd_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ex_left, y, ex_width, Inches(0.55),
        )
        cmd_box.fill.solid()
        cmd_box.fill.fore_color.rgb = RGBColor(0x2D, 0x37, 0x48)
        cmd_box.line.fill.background()
        cmd_box.adjustments[0] = 0.12
        cmd_tf = cmd_box.text_frame
        cmd_tf.word_wrap = True
        cmd_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        cmd_run = cmd_tf.paragraphs[0].add_run()
        cmd_run.text = "  " + cmd
        set_font(cmd_run, size=12, color=RGBColor(0xA0, 0xC4, 0xE8), font_family="Courier New")
        bodyPr = cmd_tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
            bodyPr.set('lIns', str(Emu(Inches(0.15))))

    # Blue accent badge
    add_rounded_rect(slide, Inches(0.9), Inches(4.5), Inches(3.5), Inches(0.5),
                     MEDIUM_BLUE, "~5-10K tokens  |  $0.01-0.05 per task", font_size=15, font_color=WHITE)


def add_token_efficiency_slide(slide):
    """Slide 6: Token Efficiency with table."""
    add_header_bar(slide, "Token Efficiency")
    add_footer_bar(slide)

    # Table
    rows, cols = 4, 3
    tbl_left = Inches(1.5)
    tbl_top = Inches(1.6)
    tbl_width = Inches(10)
    tbl_height = Inches(2.5)

    table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(2.5)

    # Data
    headers = ["Task Type", "Tokens Used", "Cost"]
    data = [
        ["Simple (calendar / tasks / email)", "0", "Free"],
        ["Complex (analysis / drafting)", "~5-10K", "~$0.01-0.05"],
        ["Idle", "0", "Free"],
    ]

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        set_font(run, size=15, bold=True, color=WHITE)
        # Background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = TABLE_HEADER_BG
        # Margins
        cell.margin_left = Inches(0.15)
        cell.margin_right = Inches(0.15)
        cell.margin_top = Inches(0.08)
        cell.margin_bottom = Inches(0.08)

    # Style data rows
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            color = ACCENT_GREEN if val == "Free" or val == "0" else BLACK
            set_font(run, size=14, color=color, bold=(val == "Free"))
            # Alternating row color
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = TABLE_ALT_ROW if i % 2 == 0 else WHITE
            cell.margin_left = Inches(0.15)
            cell.margin_right = Inches(0.15)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)

    # Additional bullet points below table
    extra_bullets = [
        '"Save" keyword also persists summaries to Obsidian -- built-in memory',
        "Previous approach used tokens for ALL tasks including simple ones",
    ]

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(extra_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        rb = p.add_run()
        rb.text = "\u2022  "
        set_font(rb, size=15, color=MEDIUM_BLUE, bold=True)
        rt = p.add_run()
        rt.text = bullet
        set_font(rt, size=15, color=BLACK)


def add_security_slide(slide):
    """Slide 7: Security."""
    add_header_bar(slide, "Security")
    add_footer_bar(slide)

    # Three protection layers as styled cards
    layers = [
        ("1", "Lark Verification Token", "Rejects all non-Lark requests"),
        ("2", "Sender Filter", "Only processes messages from authorized user (open_id)"),
        ("3", "Chat Filter", "Only processes messages from the Claude Channel (chat_id)"),
    ]

    card_width = Inches(3.5)
    card_height = Inches(2.5)
    gap = Inches(0.4)
    start_x = Inches(0.9)
    card_top = Inches(1.5)

    for i, (num, title, desc) in enumerate(layers):
        x = start_x + i * (card_width + gap)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, card_top, card_width, card_height,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BLUE
        card.line.color.rgb = MEDIUM_BLUE
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.06

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.2), card_top + Inches(0.2),
            Inches(0.55), Inches(0.55),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = DARK_BLUE
        circle.line.fill.background()
        circle_tf = circle.text_frame
        circle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle_run = circle_tf.paragraphs[0].add_run()
        circle_run.text = num
        set_font(circle_run, size=18, bold=True, color=WHITE)
        bodyPr = circle_tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')

        # Title
        add_textbox(slide, x + Inches(0.2), card_top + Inches(0.9),
                    card_width - Inches(0.4), Inches(0.4),
                    title, size=16, bold=True, color=DARK_BLUE)

        # Description
        add_textbox(slide, x + Inches(0.2), card_top + Inches(1.4),
                    card_width - Inches(0.4), Inches(0.8),
                    desc, size=13, color=GRAY)

    # Additional security notes
    extra = [
        "ngrok tunnel encrypted with TLS",
        "No API keys stored in code (uses session auth)",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(extra):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        rb = p.add_run()
        rb.text = "\u2022  "
        set_font(rb, size=15, color=ACCENT_GREEN, bold=True)
        rt = p.add_run()
        rt.text = bullet
        set_font(rt, size=15, color=BLACK)


def add_setup_steps_1_slide(slide):
    """Slide 8: Setup Steps (1/2)."""
    add_header_bar(slide, "Setup Steps (1/2)")
    add_footer_bar(slide)

    steps = [
        "Install prerequisites: lark-cli, Bun runtime, ngrok, Claude Code",
        "Create a Lark bot app on Lark Open Platform",
        "Authenticate lark-cli with read-only scopes",
        'Create "Claude Channel" group in Lark (bot + user)',
        "Build the unified channel server (TypeScript / Bun)",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)

        # Step number
        num_run = p.add_run()
        num_run.text = f"  {i + 1}.   "
        set_font(num_run, size=18, bold=True, color=MEDIUM_BLUE)

        # Step text
        step_run = p.add_run()
        step_run.text = step
        set_font(step_run, size=17, color=BLACK)


def add_setup_steps_2_slide(slide):
    """Slide 9: Setup Steps (2/2)."""
    add_header_bar(slide, "Setup Steps (2/2)")
    add_footer_bar(slide)

    steps = [
        ("6", "Configure ngrok tunnel \u2192 localhost:8765"),
        ("7", "Set Lark Event Subscription URL to ngrok URL + /webhook"),
        ("8", "Add event: im.message.receive_v1 (Message received)"),
        ("9", "Start Claude Code with: claude --dangerously-load-development-channels"),
        ("10", 'Test: send "show my calendar" in the Claude Channel'),
    ]

    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (num, step) in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)

        num_run = p.add_run()
        num_run.text = f"  {num}.   "
        set_font(num_run, size=18, bold=True, color=MEDIUM_BLUE)

        step_run = p.add_run()
        step_run.text = step
        set_font(step_run, size=17, color=BLACK)


def add_evolution_slide(slide):
    """Slide 10: Architecture Evolution."""
    add_header_bar(slide, "Architecture Evolution")
    add_footer_bar(slide)

    versions = [
        ("v1", "Cron Polling", "15s latency, tokens on every check", RGBColor(0xE5, 0x3E, 0x3E)),
        ("v2", "Background Daemon + Inbox File", "Zero idle cost, but still polling", ACCENT_ORANGE),
        ("v3", "Webhook + claude -p", "Instant simple tasks, but 70s cold start for complex", RGBColor(0xD6, 0x9E, 0x2E)),
        ("v4", "Unified Channel Server", "Instant everything, warm session, zero idle cost", ACCENT_GREEN),
    ]

    start_y = Inches(1.5)
    row_height = Inches(1.15)
    label_x = Inches(0.9)
    bar_x = Inches(2.8)
    bar_width = Inches(9.5)

    for i, (ver, name, desc, color) in enumerate(versions):
        y = start_y + i * row_height

        # Version badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            label_x, y, Inches(1.0), Inches(0.5),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        badge.adjustments[0] = 0.2
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_run = badge_tf.paragraphs[0].add_run()
        badge_run.text = ver
        set_font(badge_run, size=14, bold=True, color=WHITE)
        bodyPr = badge_tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')

        # Arrow connecting badges (except last)
        if i < len(versions) - 1:
            add_down_arrow(slide, label_x + Inches(0.4), y + Inches(0.55),
                          Pt(16), Inches(0.35), color=GRAY)

        # Name
        add_textbox(slide, bar_x, y - Inches(0.05), Inches(5), Inches(0.4),
                    name, size=16, bold=True, color=DARK_BLUE)

        # Description
        add_textbox(slide, bar_x, y + Inches(0.35), Inches(8), Inches(0.35),
                    desc, size=13, color=GRAY, italic=True)

    # Key lesson
    lesson_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.6),
    )
    lesson_box.fill.solid()
    lesson_box.fill.fore_color.rgb = LIGHT_BLUE
    lesson_box.line.color.rgb = MEDIUM_BLUE
    lesson_box.line.width = Pt(1)
    lesson_box.adjustments[0] = 0.15
    lesson_tf = lesson_box.text_frame
    lesson_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    lesson_run = lesson_tf.paragraphs[0].add_run()
    lesson_run.text = "Key lesson: each iteration optimized for reliability + latency + token efficiency"
    set_font(lesson_run, size=15, bold=True, color=DARK_BLUE)
    bodyPr = lesson_tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def add_whats_next_slide(slide):
    """Slide 11: What's Next."""
    add_header_bar(slide, "What's Next")
    add_footer_bar(slide)

    bullets = [
        "Claude Code Channels is in research preview -- will stabilize",
        "Potential to extend: add more Lark event types (doc updates, approval flows)",
        "Multi-user support: each leader gets their own Claude Channel",
        "Obsidian vault as shared team knowledge base",
        "Vision: Every ADVANCE.AI leader has an AI assistant in their pocket via Lark",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)

        rb = p.add_run()
        rb.text = "\u2022  "
        set_font(rb, size=17, color=MEDIUM_BLUE, bold=True)

        rt = p.add_run()
        rt.text = bullet
        set_font(rt, size=17, color=BLACK)

        # Highlight the last bullet (vision)
        if i == len(bullets) - 1:
            set_font(rt, size=17, color=DARK_BLUE, bold=True)


def add_thank_you_slide(slide):
    """Slide 12: Thank You."""
    # Full dark blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.4), Inches(2.2),
        Inches(2.5), Pt(4),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GREEN
    accent.line.fill.background()

    # Thank You text
    add_textbox(
        slide, Inches(1), Inches(2.5),
        Inches(11.3), Inches(1),
        "Thank You",
        size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
        font_family=FONT_HEADING,
    )

    # Built with line
    add_textbox(
        slide, Inches(1), Inches(3.7),
        Inches(11.3), Inches(0.5),
        "Built with Claude Code + Lark CLI + Obsidian",
        size=16, color=RGBColor(0xA0, 0xC4, 0xE8), alignment=PP_ALIGN.CENTER,
    )

    # Questions?
    add_textbox(
        slide, Inches(1), Inches(4.5),
        Inches(11.3), Inches(0.6),
        "Questions?",
        size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
        font_family=FONT_HEADING,
    )

    # Contact info
    add_textbox(
        slide, Inches(1), Inches(5.5),
        Inches(11.3), Inches(0.4),
        "Dong Shou", size=15, bold=True,
        color=RGBColor(0xCB, 0xD5, 0xE0), alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, Inches(1), Inches(5.9),
        Inches(11.3), Inches(0.4),
        "shoudong@advancegroup.com", size=13,
        color=RGBColor(0x90, 0xA0, 0xB0), alignment=PP_ALIGN.CENTER,
    )


if __name__ == "__main__":
    create_presentation()
